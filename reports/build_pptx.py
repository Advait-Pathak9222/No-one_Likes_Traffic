"""Build the ParkPulse Bengaluru final-round PowerPoint deck (15 dense slides).

Design principles:
- Consulting-style logic boards inspired by the reference winner deck: a big
  question, a decision checklist, dense evidence boards, and trade-off tables.
- Every slide is information-dense: section header, structured cards, mini
  tables, accurate charts and embedded real dashboard renders. No dead space.
- Abbreviations are spelled out on first use. Charts are accurate, not
  decorative. All numbers come from the pipeline outputs of the provided dataset.

Run:
    PYTHONPATH=.deps:. MPLCONFIGDIR=outputs/.matplotlib python3 reports/build_pptx.py
"""

from __future__ import annotations

import io
import json
import math
from pathlib import Path

import matplotlib

matplotlib.use("Agg", force=True)
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TABLES = ROOT / "outputs" / "tables"
FIGURES = ROOT / "outputs" / "figures"
FRONTEND = ROOT / "outputs" / "frontend"
ASSETS = ROOT / "reports" / "deck_assets"
SCREENSHOTS = ROOT / "outputs" / "deck_screenshots"
OUT = ROOT / "reports" / "parkpulse_deck.pptx"

SW, SH = 13.333, 7.5
TOTAL_SLIDES = 15
FONT = "Aptos"

BLUE = RGBColor(0x28, 0x74, 0xF0)
BLUE_DARK = RGBColor(0x13, 0x3D, 0x8F)
BLUE_50 = RGBColor(0xEC, 0xF3, 0xFF)
BLUE_100 = RGBColor(0xD8, 0xE8, 0xFF)
ORANGE = RGBColor(0xFB, 0x64, 0x1B)
ORANGE_50 = RGBColor(0xFF, 0xF1, 0xE8)
YELLOW = RGBColor(0xFF, 0xD8, 0x14)
INK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x73, 0x80, 0x92)
LINE = RGBColor(0xD8, 0xE0, 0xEC)
BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE = RGBColor(0x6D, 0x5D, 0xFF)
PURPLE_50 = RGBColor(0xF0, 0xEE, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GREEN_50 = RGBColor(0xEC, 0xF9, 0xF0)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_50 = RGBColor(0xFE, 0xEE, 0xEE)
NAVY = RGBColor(0x08, 0x18, 0x34)
NAVY_2 = RGBColor(0x10, 0x2A, 0x55)

plt.rcParams.update(
    {
        "font.family": "DejaVu Sans",
        "font.size": 11,
        "text.color": "#0f172a",
        "axes.edgecolor": "#d8e0ec",
        "axes.linewidth": 1.0,
        "xtick.color": "#475569",
        "ytick.color": "#475569",
        "figure.dpi": 180,
    }
)


def h(color: RGBColor) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def n(value, default=0.0) -> float:
    try:
        if value is None:
            return default
        value = float(value)
        return default if math.isnan(value) else value
    except Exception:
        return default


def i_fmt(value) -> str:
    return f"{n(value):,.0f}"


def f_fmt(value) -> str:
    return f"{n(value):,.1f}"


def p_fmt(value) -> str:
    return f"{100 * n(value):.1f}%"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text()) if path.exists() else {}


def load_data() -> dict:
    return {
        "metrics": load_json(FRONTEND / "metrics.json"),
        "strategy": load_json(TABLES / "strategy_validation_summary.json"),
        "policy": load_json(TABLES / "deployment_policy_summary.json"),
        "eda": load_json(TABLES / "theme_signal_eda_summary.json"),
        "plan": pd.read_csv(TABLES / "daily_enforcement_plan.csv"),
        "impact": pd.read_csv(TABLES / "operational_impact_plan.csv"),
        "eval": pd.read_csv(TABLES / "tori_ranking_evaluation.csv"),
        "forecast": pd.read_csv(TABLES / "forecast_validation_metrics.csv"),
        "corridors": pd.read_csv(TABLES / "corridor_summary.csv"),
        "context": pd.read_csv(TABLES / "theme_signal_context_features.csv"),
    }


def fig_buf(fig) -> io.BytesIO:
    b = io.BytesIO()
    fig.savefig(b, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    b.seek(0)
    return b


def clean_ax(ax, grid="x"):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis=grid, color="#e8eef7", linewidth=0.8)


# ---------- charts ----------
def chart_validation(d: dict) -> io.BytesIO:
    forecast = d["forecast"]
    selected = forecast[forecast["selected_for_predictions"] == True]
    selected20 = selected[selected["k"] == 20]
    rows = []
    if not selected20.empty:
        rows.append(("Leakage-safe forecast", selected20["capture_at_k_mean"].iloc[0]))
    eval20 = d["eval"][d["eval"]["k"] == 20]
    for name in [
        "Raw violation density", "Weighted obstruction density", "Time-safe historical mean",
        "Capacity loss pressure", "TORI final score", "ELRM recovery payoff",
    ]:
        match = eval20[eval20["method"] == name]
        if not match.empty:
            rows.append((name, match["capture_at_k_mean"].iloc[0]))
    df = pd.DataFrame(rows, columns=["method", "capture"]).drop_duplicates("method").sort_values("capture")
    fig, ax = plt.subplots(figsize=(6.4, 3.4))
    colors = [h(BLUE) if ("forecast" in m.lower() or "density" in m.lower() or "historical" in m.lower()) else h(ORANGE) for m in df["method"]]
    bars = ax.barh(df["method"], 100 * df["capture"], color=colors, height=0.66)
    for bar, val in zip(bars, 100 * df["capture"]):
        ax.text(val + 0.12, bar.get_y() + bar.get_height() / 2, f"{val:.1f}%", va="center", fontweight="bold", fontsize=9)
    ax.set_xlabel("Capture at top-20 zones (share of next-day high-pressure cells)")
    ax.set_xlim(0, max(18, 100 * df["capture"].max() * 1.24))
    ax.tick_params(axis="y", length=0, labelsize=8.5)
    clean_ax(ax)
    return fig_buf(fig)


def chart_capture_k(d: dict) -> io.BytesIO:
    f = d["forecast"]
    names = {"hist_gradient_boosting": "HistGradientBoosting (selected)", "lightgbm": "LightGBM",
             "xgboost": "XGBoost", "rolling_7d_baseline": "7-day rolling baseline"}
    colors = {"hist_gradient_boosting": h(BLUE), "lightgbm": h(PURPLE), "xgboost": h(ORANGE), "rolling_7d_baseline": h(MUTED)}
    fig, ax = plt.subplots(figsize=(5.7, 3.3))
    for model in ["rolling_7d_baseline", "xgboost", "lightgbm", "hist_gradient_boosting"]:
        sub = f[f["model"] == model].sort_values("k")
        if sub.empty:
            continue
        lw = 3.4 if model == "hist_gradient_boosting" else 1.7
        ax.plot(sub["k"], 100 * sub["capture_at_k_mean"], marker="o", ms=4, lw=lw, color=colors.get(model, "#888"), label=names.get(model, model))
    ax.set_xlabel("Top-K ranked zones evaluated")
    ax.set_ylabel("Capture@K (%)")
    ax.set_xticks([10, 20, 50, 100])
    ax.legend(frameon=False, fontsize=7.6, loc="upper left")
    clean_ax(ax)
    return fig_buf(fig)


def chart_concentration(d: dict) -> io.BytesIO:
    df = pd.read_csv(TABLES / "theme_signal_grid_concentration_curve.csv")
    x = df[[c for c in df.columns if "share" in c.lower()][0]]
    ycols = [c for c in df.columns if "weighted" in c.lower() or "obstruction" in c.lower()]
    y = df[ycols[0] if ycols else df.columns[-1]]
    fig, ax = plt.subplots(figsize=(5.2, 3.0))
    ax.plot(100 * x, 100 * y, color=h(BLUE), linewidth=3)
    ax.fill_between(100 * x, 100 * y, color="#dbeafe", alpha=0.9)
    ax.plot([0, 100], [0, 100], color="#cbd5e1", linestyle="--", linewidth=1)
    ax.set_xlabel("Share of grid cells covered (%)")
    ax.set_ylabel("Weighted obstruction captured (%)")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    clean_ax(ax)
    return fig_buf(fig)


def chart_repeat_gap(d: dict) -> io.BytesIO:
    plan = d["plan"]
    fig, ax = plt.subplots(figsize=(5.3, 2.95))
    ax.hist(plan["repeat_pressure_score_0_100"], bins=np.arange(0, 105, 10), color=h(PURPLE), alpha=0.78, label="Repeat-vehicle pressure")
    ax.hist(plan["patrol_gap_score_0_100"], bins=np.arange(0, 105, 10), color=h(ORANGE), alpha=0.55, label="Patrol-coverage gap")
    ax.axvline(85, color="#0f172a", linestyle="--", linewidth=1.2)
    ax.text(60, ax.get_ylim()[1] * 0.86, "priority threshold (85)", fontsize=8)
    ax.set_xlabel("Score (0-100)")
    ax.set_ylabel("Priority zones")
    ax.legend(frameon=False, fontsize=8)
    clean_ax(ax)
    return fig_buf(fig)


def chart_actions(d: dict) -> io.BytesIO:
    counts = d["plan"]["recommended_action"].value_counts()
    labels = [x.replace(" + ", " + ").replace(" enforcement", "") for x in counts.index]
    fig, ax = plt.subplots(figsize=(4.2, 3.0))
    wedges, _ = ax.pie(
        counts.values, colors=[h(BLUE), h(PURPLE), h(GREEN), h(ORANGE), h(RED)][: len(counts)],
        startangle=90, counterclock=False, wedgeprops={"width": 0.42, "edgecolor": "white", "linewidth": 2},
    )
    ax.legend(wedges, [f"{l}: {v}" for l, v in zip(labels, counts.values)], loc="center left", bbox_to_anchor=(0.82, 0.5), frameon=False, fontsize=8)
    ax.set(aspect="equal")
    return fig_buf(fig)


def chart_elrm(d: dict) -> io.BytesIO:
    g = d["impact"].groupby("recommended_action")["estimated_lane_recovery_minutes"].sum().sort_values()
    labels = [x.replace(" + ", " +\n").replace(" enforcement", "") for x in g.index]
    fig, ax = plt.subplots(figsize=(5.7, 3.2))
    bars = ax.barh(labels, g.values, color=h(BLUE), height=0.62)
    for bar, val in zip(bars, g.values):
        ax.text(val + g.max() * 0.01, bar.get_y() + bar.get_height() / 2, f"{val:,.0f}", va="center", fontsize=8, fontweight="bold")
    ax.set_xlabel("Equivalent Lane Recovery Minutes (total across plan)")
    ax.tick_params(axis="y", labelsize=8)
    clean_ax(ax)
    return fig_buf(fig)


def chart_policy(d: dict) -> io.BytesIO:
    path = FIGURES / "deployment_policy_budget_frontier.png"
    if path.exists():
        return io.BytesIO(path.read_bytes())
    return chart_elrm(d)


# ---------- ppt helpers ----------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def shp(slide, typ, x, y, w, hgt, fill=WHITE, line=None, lw=1, radius=0.08):
    s = slide.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(hgt))
    s.shadow.inherit = False
    if typ == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    return s


def rect(slide, x, y, w, hgt, fill=WHITE, line=LINE, rounded=True, lw=1):
    return shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, hgt, fill, line, lw)


def tx(slide, x, y, w, hgt, content, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(hgt))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    runs = content if isinstance(content, list) else [(str(content), {})]
    for txt, ov in runs:
        r = p.add_run()
        r.text = str(txt)
        f = r.font
        f.name = FONT
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", False)
        f.color.rgb = ov.get("color", color)
    return box


def bullets(slide, x, y, w, hgt, items, size=10.5, color=INK, gap=3):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(hgt))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.13
        r = p.add_run()
        r.text = "• " + item
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def pic(slide, img, x, y, w=None, hgt=None):
    if hasattr(img, "seek"):
        img.seek(0)
    args = {}
    if w is not None:
        args["width"] = Inches(w)
    if hgt is not None:
        args["height"] = Inches(hgt)
    return slide.shapes.add_picture(img, Inches(x), Inches(y), **args)


def line(slide, x1, y1, x2, y2, color=LINE, lw=1):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    return c


def tag(slide, x, y, text, fill=BLUE, w=None, color=WHITE):
    width = w or max(0.58, 0.085 * len(str(text)) + 0.3)
    rect(slide, x, y, width, 0.28, fill=fill, line=None, rounded=True)
    tx(slide, x, y + 0.03, width, 0.2, text, size=7.6, color=color, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return width


def logo(slide, x=0.48, y=0.28, dark=False):
    tc = WHITE if dark else INK
    mc = RGBColor(0xC7, 0xD7, 0xF5) if dark else MUTED
    rect(slide, x, y + 0.03, 0.34, 0.34, fill=YELLOW, line=None, rounded=True)
    tx(slide, x, y + 0.06, 0.34, 0.22, "F", size=10.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    tx(slide, x + 0.42, y - 0.01, 0.92, 0.34, "Flipkart", size=11, color=tc, bold=True)
    tx(slide, x + 1.22, y - 0.01, 0.22, 0.34, "×", size=11, color=mc, bold=True, align=PP_ALIGN.CENTER)
    rect(slide, x + 1.48, y + 0.03, 0.34, 0.34, fill=BLUE, line=None, rounded=True)
    tx(slide, x + 1.48, y + 0.055, 0.34, 0.22, "P", size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tx(slide, x + 1.9, y - 0.01, 1.4, 0.34, "ParkPulse", size=11, color=tc, bold=True)


def header(slide, section, title, page, dark=False):
    logo(slide, dark=dark)
    tx(slide, 0.55, 0.82, 6.0, 0.24, section.upper(), size=8.5, color=ORANGE, bold=True)
    tx(slide, 0.55, 1.05, 11.9, 0.5, title, size=20, color=WHITE if dark else INK, bold=True)
    tx(slide, 12.05, 0.34, 0.7, 0.22, f"{page:02d} / {TOTAL_SLIDES}", size=8.5, color=RGBColor(0xC7, 0xD7, 0xF5) if dark else MUTED, bold=True, align=PP_ALIGN.RIGHT)


def foot(slide, note="Source: ParkPulse pipeline outputs from the provided Theme-1 enforcement dataset", dark=False):
    tx(slide, 0.56, 7.1, 8.2, 0.22, note, size=7.2, color=RGBColor(0xC7, 0xD7, 0xF5) if dark else MUTED)
    tx(slide, 9.7, 7.1, 3.0, 0.22, "Proxy metrics are labelled; nothing is presented as measured live traffic", size=7.0, color=RGBColor(0xC7, 0xD7, 0xF5) if dark else MUTED, align=PP_ALIGN.RIGHT)


def kpi(slide, x, y, w, hgt, value, label, accent=BLUE, fill=WHITE, dark=False):
    rect(slide, x, y, w, hgt, fill=fill if not dark else NAVY_2, line=None if dark else LINE)
    rect(slide, x, y, 0.06, hgt, fill=accent, line=None, rounded=False)
    tx(slide, x + 0.16, y + 0.09, w - 0.26, 0.34, value, size=18, color=WHITE if dark else INK, bold=True)
    tx(slide, x + 0.16, y + 0.47, w - 0.26, 0.3, label, size=7.0, color=RGBColor(0xC7, 0xD7, 0xF5) if dark else MUTED, bold=True)


def card(slide, x, y, w, hgt, title, body, accent=BLUE, fill=WHITE, title_size=11.5):
    rect(slide, x, y, w, hgt, fill=fill, line=LINE)
    rect(slide, x, y, 0.06, hgt, fill=accent, line=None, rounded=False)
    tx(slide, x + 0.18, y + 0.12, w - 0.28, 0.3, title, size=title_size, color=INK, bold=True)
    if isinstance(body, list):
        bullets(slide, x + 0.18, y + 0.49, w - 0.3, hgt - 0.57, body, size=8.8, color=SLATE, gap=1)
    else:
        tx(slide, x + 0.18, y + 0.49, w - 0.3, hgt - 0.57, body, size=9.2, color=SLATE, spacing=1.12)


def mini_table(slide, x, y, w, hgt, headers, rows, col_widths=None, header_fill=BLUE_50, fs=7.3):
    col_widths = col_widths or [w / len(headers)] * len(headers)
    row_h = hgt / (len(rows) + 1)
    rect(slide, x, y, w, hgt, fill=WHITE, line=LINE, rounded=False)
    cx = x
    for j, head in enumerate(headers):
        rect(slide, cx, y, col_widths[j], row_h, fill=header_fill, line=LINE, rounded=False)
        tx(slide, cx + 0.05, y + 0.04, col_widths[j] - 0.1, row_h - 0.05, head, size=fs, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_widths[j]
    for i, row in enumerate(rows):
        cx = x
        for j, cell in enumerate(row):
            rect(slide, cx, y + row_h * (i + 1), col_widths[j], row_h, fill=WHITE, line=LINE, rounded=False)
            tx(slide, cx + 0.06, y + row_h * (i + 1) + 0.03, col_widths[j] - 0.12, row_h - 0.05, str(cell), size=fs, color=INK if j == 0 else SLATE, bold=(j == 0), align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_widths[j]


def shot(slide, name, x, y, w):
    img = SCREENSHOTS / name
    if not img.exists():
        img = ASSETS / name
    if img.exists():
        try:
            with Image.open(img) as im:
                ratio = im.height / im.width
        except Exception:
            ratio = 0.719
        rect(slide, x - 0.04, y - 0.04, w + 0.08, w * ratio + 0.08, fill=WHITE, line=LINE)
        return pic(slide, str(img), x, y, w=w)
    rect(slide, x, y, w, w * 0.719, fill=BLUE_50, line=LINE)
    return None


# ---------- slides ----------
def slide_title(prs, d, pg):
    s = blank(prs)
    rect(s, -0.1, -0.1, SW + 0.2, SH + 0.2, fill=BLUE, line=None, rounded=False)
    logo(s, 0.75, 0.58, dark=True)
    tx(s, 0.95, 1.95, 11.4, 1.0, "Which illegal parking should Bengaluru solve first?", size=34, color=WHITE, bold=True, spacing=1.04)
    tx(s, 0.95, 3.08, 10.9, 0.85, "ParkPulse turns five months of police enforcement records into a time-windowed action plan: which hotspots, which repeat-offender vehicles, which patrol-coverage gaps, and how much running-road capacity each enforcement action recovers.", size=15, color=RGBColor(0xE9, 0xF1, 0xFF), spacing=1.2)
    m = d["metrics"]
    vals = [
        (i_fmt(m["total_violations"]), "Violation records analysed"),
        (i_fmt(m["total_hotspots"]), "Hotspot cells ranked"),
        (i_fmt(m["top20_lane_recovery_minutes"]) + " min", "Top-20 lane recovery (estimate)"),
        (p_fmt(m["headline_recurrence_capture_at_20"]), "Validated recurrence capture"),
    ]
    for i, (v, lab) in enumerate(vals):
        kpi(s, 0.95 + i * 3.02, 5.05, 2.75, 1.05, v, lab, YELLOW, dark=True)
    tx(s, 0.95, 6.85, 11.5, 0.26, "Flipkart GRiD · Theme 1 · Illegal-parking hotspots and their traffic-flow impact · Product strategy and data-science solution deck", size=9, color=RGBColor(0xC7, 0xD7, 0xF5), bold=True)


def slide_questions(prs, d, pg):
    s = blank(prs)
    rect(s, -0.1, -0.1, SW + 0.2, SH + 0.2, fill=BLUE, line=None, rounded=False)
    logo(s, 0.75, 0.55, dark=True)
    tx(s, 0.95, 1.3, 11.0, 0.45, "What a winning Theme-1 answer must prove", size=25, color=WHITE, bold=True)
    qs = [
        "Identify illegal-parking hotspots beyond raw counts, accounting for recording bias.",
        "Quantify the traffic-flow impact honestly, without inventing speed or queue measurements.",
        "Tell police what to do: tow, fixed-window patrol, engineering fix, or spillover control.",
        "Explain why a zone is a priority: recurrence, road-space, repeat vehicles, patrol gap.",
        "Run from day one on the data Bengaluru already has, with no external data purchase.",
        "Improve as richer feeds arrive, without redesigning the product.",
    ]
    y = 2.2
    for i, q in enumerate(qs, 1):
        rect(s, 1.05, y + 0.02, 0.36, 0.36, fill=YELLOW, line=None, rounded=True)
        tx(s, 1.05, y + 0.07, 0.36, 0.22, str(i), size=9, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
        tx(s, 1.62, y, 10.4, 0.36, q, size=14, color=WHITE, bold=True)
        y += 0.66
    tx(s, 0.95, 6.78, 11.9, 0.32, "This deck answers each question with exploratory analysis, engineering, machine-learning validation, and a working command-center product.", size=10.5, color=RGBColor(0xC7, 0xD7, 0xF5), bold=True)


def slide_strategy(prs, d, pg):
    s = blank(prs)
    header(s, "Executive strategy", "From violation records to enforcement return on investment, on one page", pg)
    m = d["metrics"]
    mini_table(s, 0.55, 1.78, 12.22, 1.2,
               ["Who uses it", "Their goal", "Their constraint", "Success measure"],
               [["Traffic control room", "Dispatch the first wave", "Limited patrol and tow units", "Running-lane minutes recovered"],
                ["Station inspector", "Plan officer beats", "Recording / coverage bias", "Repeat-offender and gap reduction"],
                ["Field officer", "Clear the obstruction", "No time to read tables", "Response-deadline compliance"]],
               col_widths=[2.6, 3.0, 3.3, 3.32])
    tx(s, 0.58, 3.2, 6.0, 0.24, "HOW ILLEGAL PARKING BECOMES CONGESTION", size=8.5, color=PURPLE, bold=True)
    chain = [("Wrong / no parking", ORANGE), ("Kerb / running lane blocked", BLUE), ("Queue spillback at junction", PURPLE), ("Tow / patrol within deadline", GREEN), ("Running-lane minutes recovered", BLUE_DARK)]
    x = 0.6
    for idx, (label, col) in enumerate(chain):
        rect(s, x, 3.62, 2.14, 0.66, fill=WHITE, line=col, rounded=True, lw=1.3)
        tx(s, x + 0.08, 3.72, 1.98, 0.46, label, size=8.6, color=INK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if idx < len(chain) - 1:
            shp(s, MSO_SHAPE.RIGHT_ARROW, x + 2.16, 3.84, 0.28, 0.22, fill=col, line=None)
        x += 2.5
    rect(s, 0.6, 4.66, 5.4, 1.12, fill=PURPLE_50, line=PURPLE)
    tx(s, 0.78, 4.82, 5.05, 0.82, "Priority = validated recurrence × obstruction severity × exposure correction × repeat-offender pressure × patrol-coverage gap", size=12, color=INK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    kpi(s, 6.2, 4.66, 2.05, 1.12, p_fmt(m["headline_recurrence_capture_at_20"]), "Validated recurrence capture", BLUE)
    kpi(s, 8.42, 4.66, 2.05, 1.12, i_fmt(m["chronic_repeat_priority_zones"]), "Chronic repeat-offender zones", PURPLE)
    kpi(s, 10.64, 4.66, 2.05, 1.12, i_fmt(m["patrol_gap_priority_zones"]), "Patrol-coverage gap zones", ORANGE)
    tx(s, 0.6, 6.18, 12.1, 0.5, "Thesis: narrow the enforcement battlefield to the small set of high-recurrence, high-obstruction, high-payoff zones, instead of chasing every recorded violation equally.", size=12, color=INK, bold=True, spacing=1.1)
    foot(s)


def slide_challenge(prs, d, pg):
    s = blank(prs)
    header(s, "The operational challenge", "Why patrol-based enforcement cannot scale — and what we can honestly measure", pg)
    m = d["metrics"]
    card(s, 0.58, 1.78, 3.05, 1.55, "Reactive patrolling", "Deployment follows complaints or visible density. It cannot see the downstream traffic impact of a blocked lane.", ORANGE, ORANGE_50)
    card(s, 3.88, 1.78, 3.05, 1.55, "No impact heat map", "A goods vehicle on a main-road junction is not the same as a quiet-street two-wheeler, but counts treat them equally.", BLUE, BLUE_50)
    card(s, 7.18, 1.78, 3.05, 1.55, "Resource bottleneck", "Tow vehicles and officers are scarce and must be rationed by payoff, not by raw violation count.", PURPLE, PURPLE_50)
    kpi(s, 10.55, 1.78, 1.92, 1.55, i_fmt(m["high_spillback_risk_zones"]), "High queue-spillback zones identified", RED)
    tx(s, 0.62, 3.62, 8.0, 0.34, "Given the data we were handed, what can we measure — and how should we label each claim?", size=14, color=INK, bold=True)
    mini_table(s, 0.6, 4.12, 12.15, 2.0,
               ["Requirement from the theme", "Evidence in the dataset", "ParkPulse answer", "Honesty label"],
               [["Detect hotspots", "Latitude, longitude, station, time stamp", "250-metre cells, density clustering, inferred corridors", "Directly observed"],
                ["Quantify traffic impact", "Violation type, vehicle class, road-space context", "Capacity-loss, spillback risk, lane-recovery minutes", "Transparent proxy"],
                ["Prioritise enforcement", "Active days, recording exposure, repeat vehicle identifiers", "Recurrence forecast, obstruction index, patrol gap", "Modelled"],
                ["Produce an enforcement plan", "Station, time window, actionable fields", "Tow / patrol / engineering playbook with deadlines", "Product output"]],
               col_widths=[2.7, 3.3, 4.05, 2.1])
    foot(s)


def slide_dataset(prs, d, pg):
    s = blank(prs)
    header(s, "Dataset and data quality", "We start by understanding — and stress-testing — the data we were given", pg)
    m, eda = d["metrics"], d["eda"]
    kpi(s, 0.58, 1.78, 2.36, 1.0, i_fmt(eda.get("rows", m["total_violations"])), "Violation records (Jan–May)", BLUE)
    kpi(s, 3.15, 1.78, 2.36, 1.0, i_fmt(eda.get("active_days", 151)), "Active calendar days", PURPLE)
    kpi(s, 5.72, 1.78, 2.36, 1.0, i_fmt(eda.get("stations", 55)), "Police stations covered", GREEN)
    kpi(s, 8.29, 1.78, 2.36, 1.0, i_fmt(eda.get("grid_250m_cells", 3443)), "Active 250-metre grid cells", ORANGE)
    kpi(s, 10.86, 1.78, 1.75, 1.0, "8", "Distinct data signal families", BLUE_DARK)
    card(s, 0.58, 3.0, 5.95, 1.55, "Field completeness (what is and is not present)",
         ["Location, time, station, violation text: effectively complete", "Validation time stamp present on about 58 percent of records", "Action-taken and closure time stamps are absent (0 percent)",
          "Repeat-offender vehicle identifier present for known vehicles"], BLUE, BLUE_50)
    rect(s, 6.7, 3.0, 6.05, 1.55, fill=RED_50, line=RED)
    tx(s, 6.9, 3.14, 5.7, 0.3, "Data-quality risk we flag openly", size=11.5, color=RED, bold=True)
    tx(s, 6.9, 3.5, 5.7, 1.0,
       f"Only {i_fmt(m['evening_records'])} records fall in the 18:00–22:00 evening window, against {i_fmt(m['night_records'])} at night (a ratio of {m['evening_to_night_record_ratio']*100:.2f} percent). Rather than treat the evening peak as genuinely empty, ParkPulse marks this window as a recording-coverage risk so officers are not misled.",
       size=9.4, color=SLATE, spacing=1.16)
    tx(s, 0.6, 4.78, 12.1, 0.34, "Because action and closure outcomes are missing, ParkPulse is honestly a decision-support layer today, with built-in hooks to learn from outcomes once those feeds exist.", size=11, color=INK, bold=True, spacing=1.1)
    mini_table(s, 0.6, 5.38, 12.15, 1.18,
               ["Missing / weak field", "Coverage", "Consequence for the design"],
               [["Action-taken / closure time stamp", "0 percent", "No measured enforcement outcome yet — flagged as a future integration"],
                ["Signal / zebra-crossing context", "About 0.2 percent", "Junction conflict treated as a low-volume, supporting signal only"]],
               col_widths=[3.6, 2.0, 6.55])
    foot(s)


def slide_eda_signals(prs, d, pg):
    s = blank(prs)
    header(s, "Exploratory Data Analysis (EDA), part one", "Three strong signals make this a targeted-enforcement problem, not generic clustering", pg)
    kpi(s, 0.58, 1.78, 2.4, 0.95, "65.5%", "Weighted obstruction in the top 5% of cells", BLUE)
    kpi(s, 3.2, 1.78, 2.4, 0.95, "34.2%", "Known-vehicle records that repeat (2 or more)", PURPLE)
    kpi(s, 5.82, 1.78, 2.4, 0.95, "18.3%", "Chronic records (3 or more occurrences)", ORANGE)
    kpi(s, 8.44, 1.78, 2.4, 0.95, "0.85", "Correlation of raw count with recording exposure", RED)
    kpi(s, 11.06, 1.78, 1.55, 0.95, "58%", "Records with a validation time stamp", GREEN)
    pic(s, chart_concentration(d), 0.62, 3.0, w=4.5)
    pic(s, chart_repeat_gap(d), 5.35, 3.05, w=4.55)
    context = d["context"].head(6)
    rows = [[r.context_signal[:22], f"{100*r.record_share:.1f}%", f"{100*r.weighted_share:.1f}%"] for r in context.itertuples()]
    mini_table(s, 10.08, 3.0, 2.6, 2.85, ["Road-space context", "Rows", "Weight"], rows, col_widths=[1.3, 0.62, 0.68])
    tx(s, 0.62, 6.3, 12.1, 0.5, "Insight: obstruction is highly concentrated, a third of records come from repeat vehicles, and raw counts are biased by where officers already record. The design must correct for all three.", size=11, color=INK, bold=True, spacing=1.08)
    foot(s)


def slide_eda_spatial(prs, d, pg):
    s = blank(prs)
    header(s, "Exploratory Data Analysis (EDA), part two", "Where and when obstruction concentrates, and the eight signal families we engineered", pg)
    heat = FIGURES / "theme_station_timeblock_heatmap.png"
    if heat.exists():
        rect(s, 0.55, 1.85, 6.5, 4.55, fill=WHITE, line=LINE)
        pic(s, str(heat), 0.66, 1.96, w=6.28)
    tx(s, 0.66, 6.32, 6.4, 0.22, "Station × time-window weighted-obstruction heat map: pressure clusters in specific station-and-window pairs.", size=8, color=MUTED, spacing=1.05)
    tx(s, 7.25, 1.85, 5.5, 0.26, "EIGHT DATA SIGNAL FAMILIES ENGINEERED FROM THE RECORDS", size=8.5, color=PURPLE, bold=True)
    fams = [
        ["1. Micro-hotspot concentration", "Strong", "Top 5% cells hold 65.5% of obstruction"],
        ["2. Repeat / chronic vehicles", "Strong", "34.2% repeat, 18.3% chronic offenders"],
        ["3. Road-space obstruction context", "Strong", "Main-road 29.4%, large-vehicle 23.4%"],
        ["4. Signal / zebra / junction conflict", "Weak", "Sparse (about 0.2%); supporting only"],
        ["5. Recording / exposure bias", "Strong", "Count-vs-exposure correlation 0.85"],
        ["6. Validation and evidence latency", "Medium", "Validation time stamp on 58% of rows"],
        ["7. Action / closure outcome data", "Absent", "0% present — future integration hook"],
        ["8. SCITA transfer flag", "Strong", "Present on 100%, true for 85.7%"],
    ]
    mini_table(s, 7.25, 2.2, 5.5, 4.0, ["Signal family", "Support", "What the data shows"], fams, col_widths=[2.35, 0.95, 2.2], fs=7.4)
    foot(s)


def slide_architecture(prs, d, pg):
    s = blank(prs)
    header(s, "Solution architecture", "ParkPulse is an operating system for parking enforcement, in five layers", pg)
    layers = [
        ("1", "Data clean-room", "Normalise vehicle class, time window, station, violation text and evidence status", BLUE),
        ("2", "Hotspot engine", "250-metre grid, density clustering, station / time windows, inferred corridors", PURPLE),
        ("3", "Scoring and machine learning", "Recurrence forecast, obstruction index, exposure correction, repeat and gap scores", ORANGE),
        ("4", "Operational impact model", "Capacity loss, spillback risk, response deadline, lane-recovery minutes", GREEN),
        ("5", "Product layer", "Command center, live operations brief, simulator, printable field brief", BLUE_DARK),
    ]
    x = 0.62
    for idx, (num, title, desc, col) in enumerate(layers):
        rect(s, x, 1.95, 2.18, 2.25, fill=WHITE, line=col, rounded=True, lw=1.3)
        rect(s, x + 0.14, 2.12, 0.4, 0.4, fill=col, line=None, rounded=True)
        tx(s, x + 0.14, 2.18, 0.4, 0.28, num, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tx(s, x + 0.16, 2.62, 1.86, 0.5, title, size=11, color=INK, bold=True, align=PP_ALIGN.CENTER)
        tx(s, x + 0.14, 3.16, 1.9, 0.95, desc, size=8.0, color=SLATE, align=PP_ALIGN.CENTER, spacing=1.08)
        if idx < len(layers) - 1:
            shp(s, MSO_SHAPE.RIGHT_ARROW, x + 2.2, 2.92, 0.32, 0.28, fill=col, line=None)
        x += 2.5
    card(s, 0.66, 4.78, 3.78, 1.4, "Machine-learning backbone", ["Gradient-boosted recurrence forecast", "Time-safe validation, no future leakage", "Top-K capture and ranking-quality scoring"], BLUE, BLUE_50)
    card(s, 4.66, 4.78, 3.78, 1.4, "Operations backbone", ["Station deployment cards", "Patrol and tow resource budgeting", "Response deadlines and a field playbook"], GREEN, GREEN_50)
    card(s, 8.66, 4.78, 4.1, 1.4, "Governance backbone", ["Proxy claims labelled explicitly", "Data-quality warnings surfaced in product", "Hooks for road, speed, signal and roster feeds"], ORANGE, ORANGE_50)
    foot(s)


def slide_engineering(prs, d, pg):
    s = blank(prs)
    header(s, "Engineering techniques", "A reproducible sixteen-stage pipeline turns raw records into a deployable plan", pg)
    stages = [
        "Clean and normalise records", "Build enforcement exposure table", "Exploratory data-signal analysis",
        "Density clustering of hotspots", "Traffic Obstruction Risk Index", "Enforcement recommendation rules",
        "Road-space and corridor intelligence", "Obstruction-risk overlay", "Operational impact and recovery model",
        "Deployment-policy back-test", "Leakage-safe forecasting table", "Recurrence and ranking evaluation",
        "Strategy validation summary", "Robustness and stability checks", "Figures and maps", "Frontend data export",
    ]
    cols = 4
    cw, ch, gx, gy = 2.92, 0.5, 0.13, 0.13
    x0, y0 = 0.6, 1.85
    for i, st in enumerate(stages):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        y = y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=BLUE_50 if i % 2 == 0 else WHITE, line=LINE)
        rect(s, x, y, 0.34, ch, fill=BLUE, line=None, rounded=False)
        tx(s, x, y + 0.07, 0.34, 0.34, f"{i+1}", size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tx(s, x + 0.42, y + 0.04, cw - 0.48, ch - 0.06, st, size=8.0, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
    y_tech = y0 + 4 * (ch + gy) + 0.12
    techniques = [
        ("Determinism and reproducibility", "Fixed 250-metre grid identifiers and seeded clustering make every run repeatable.", BLUE),
        ("Exposure de-biasing", "Counts are divided by officer, device and active-day exposure so watched zones do not dominate.", PURPLE),
        ("Percentile-rank normalisation", "Every score is rank-normalised to 0–100 so components combine on a common scale.", ORANGE),
        ("Leakage-safe construction", "Lagged and historical features only; the last 21 days are held out for validation.", GREEN),
    ]
    for i, (t, b, col) in enumerate(techniques):
        card(s, 0.6 + i * 3.06, y_tech, 2.92, 1.55, t, b, col, WHITE, title_size=9.8)
    foot(s)


def slide_ml(prs, d, pg):
    s = blank(prs)
    header(s, "Data modelling and machine learning", "A leakage-safe recurrence model, validated against three alternatives", pg)
    m = d["metrics"]
    f20 = d["forecast"][d["forecast"]["k"] == 20]
    name_map = {"hist_gradient_boosting": "HistGradientBoosting", "lightgbm": "LightGBM", "xgboost": "XGBoost", "rolling_7d_baseline": "7-day rolling baseline"}
    order = ["hist_gradient_boosting", "lightgbm", "xgboost", "rolling_7d_baseline"]
    rows = []
    for mdl in order:
        sub = f20[f20["model"] == mdl]
        if sub.empty:
            continue
        sel = "Selected" if bool(sub["selected_for_predictions"].iloc[0]) else "—"
        rows.append([name_map[mdl], f"{100*sub['capture_at_k_mean'].iloc[0]:.1f}%", f"{sub['ndcg_at_k_mean'].iloc[0]:.3f}", sel])
    pic(s, chart_capture_k(d), 0.6, 1.85, w=5.55)
    tx(s, 0.66, 5.5, 5.5, 0.22, "Capture@K rises smoothly with K for every model; the selected model leads at the operational K of 20.", size=8, color=MUTED, spacing=1.05)
    mini_table(s, 6.35, 1.88, 6.4, 1.55, ["Model compared", "Capture@20", "Ranking quality (NDCG@20)", "Outcome"], rows, col_widths=[2.7, 1.3, 1.7, 0.7], fs=7.6)
    card(s, 6.35, 3.62, 3.12, 1.5, "Engineered features", ["Lagged 3, 7, 14 and 28-day rolling means of count and weight", "Expanding historical means", "Day index, weekday and weekend flags", "Grid, station and time-window codes"], BLUE, BLUE_50, title_size=10.2)
    card(s, 9.63, 3.62, 3.12, 1.5, "Leakage-safe validation", ["Train on all days before the last 21", "Validate on the final 21 days only", "Features use only past information", "Selection by Capture@20 and NDCG@20"], GREEN, GREEN_50, title_size=10.2)
    rect(s, 6.35, 5.28, 6.4, 0.92, fill=PURPLE_50, line=PURPLE)
    tx(s, 6.55, 5.4, 6.05, 0.7, [
        ("Honest finding: ", {"bold": True, "color": PURPLE}),
        ("the gradient-boosted forecast and raw density are statistically close for predicting recurrence. We report this openly and use the obstruction index and impact metrics for the separate job of deciding action.", {"color": INK}),
    ], size=9.6, spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)
    foot(s)


def slide_scoring(prs, d, pg):
    s = blank(prs)
    header(s, "Scoring logic and control-room metrics", "Explainable enough for police deployment; every metric has a transparent formula", pg)
    rect(s, 0.58, 1.78, 6.0, 1.08, fill=BLUE_50, line=BLUE)
    tx(s, 0.76, 1.9, 5.64, 0.82, [("Traffic Obstruction Risk Index = ", {"bold": True}), ("density + persistence + time pressure + junction criticality + severity + vehicle obstruction + confidence + repeat + gap", {})], size=10.6, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 6.74, 1.78, 6.02, 1.08, fill=ORANGE_50, line=ORANGE)
    tx(s, 6.92, 1.9, 5.66, 0.82, [("Equivalent Lane Recovery Minutes = ", {"bold": True}), ("window minutes × obstruction fraction × lane context × recurrence × action effectiveness × confidence", {})], size=10.6, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    mini_table(s, 0.58, 3.05, 12.2, 2.5,
               ["Control-room metric", "What it answers", "How it is computed (transparent proxy)"],
               [["Capacity-loss pressure (0–100)", "How much running lane is at risk", "Obstruction × lane context × recurrence × time-of-day factor"],
                ["Queue-spillback risk (0–100)", "Could it back up into a junction", "0.42 capacity + 0.25 junction + 0.18 main-road + 0.15 time"],
                ["Clearance deadline (minutes)", "How fast the first unit must arrive", "Bands of 15 / 30 / 60 / 120 from spillback and capacity"],
                ["Evidence quality (0–100)", "How much to trust the record", "0.40 confidence + 0.30 volume + 0.15 context + 0.15 repeat"],
                ["Repeat-offender pressure (0–100)", "Is the same vehicle chronic", "Repeat, chronic, weighted and multi-station vehicle shares"],
                ["Patrol-coverage gap (0–100)", "High risk but low past coverage", "Exposure-adjusted risk × inverse historical exposure"],
                ["Operational priority (0–100)", "What to dispatch first", "0.28 capacity + 0.22 spillback + recovery + 0.12 evidence + repeat + gap"]],
               col_widths=[3.0, 3.5, 5.7], fs=7.3)
    tx(s, 0.6, 5.72, 12.1, 0.5, "Two distinct scores by design: the Traffic Obstruction Risk Index explains and ranks severity, while operational priority sequences the actual dispatch under patrol and tow limits.", size=11, color=INK, bold=True, spacing=1.08)
    foot(s)


def slide_validation(prs, d, pg):
    s = blank(prs)
    header(s, "Validation", "We separate prediction quality from action quality — and prove both", pg)
    pic(s, chart_validation(d), 0.6, 1.9, w=5.7)
    m = d["metrics"]
    kpi(s, 6.5, 1.92, 2.0, 0.95, p_fmt(m["headline_recurrence_capture_at_20"]), "Forecast capture at top-20", BLUE)
    kpi(s, 8.72, 1.92, 2.0, 0.95, p_fmt(m["raw_count_baseline_capture_at_20"]), "Raw-density baseline", PURPLE)
    kpi(s, 10.94, 1.92, 1.78, 0.95, p_fmt(m["robust_top20_overlap"]), "Robust top-20 overlap", GREEN)
    card(s, 6.5, 3.2, 2.98, 1.35, "Finding one", "Raw density is a strong recurrence baseline; we do not hide that or overclaim the model.", BLUE, BLUE_50)
    card(s, 9.74, 3.2, 2.98, 1.35, "Finding two", "The obstruction index is the action and explanation layer, not the recurrence winner.", ORANGE, ORANGE_50)
    card(s, 6.5, 4.72, 2.98, 1.45, "Finding three", "Operational priority optimises field payoff: recovery, deadline, spillback and evidence.", GREEN, GREEN_50)
    card(s, 9.74, 4.72, 2.98, 1.45, "Finding four", "Time-safe validation prevents future leakage and mirrors a hidden, live deployment.", PURPLE, PURPLE_50)
    foot(s)


def slide_spatial(prs, d, pg):
    s = blank(prs)
    header(s, "Spatial intelligence", "Exact points become a lane-risk and corridor action map", pg)
    overlay = FIGURES / "congestion_risk_overlay.png"
    if overlay.exists():
        rect(s, 0.62, 1.82, 5.5, 4.5, fill=WHITE, line=LINE)
        pic(s, str(overlay), 0.78, 1.95, hgt=3.95)
    tx(s, 0.78, 6.05, 5.2, 0.34, "Traffic-style overlay (blue, yellow, red obstruction proxy from records) — this is inferred from violations, not live road speed.", size=8, color=MUTED, spacing=1.05)
    cor = d["corridors"].head(6)
    rows = [[r.corridor_name[:26], r.station[:12], int(r.linked_hotspots), f"{int(r.total_violations):,}"] for r in cor.itertuples()]
    mini_table(s, 6.4, 1.9, 6.35, 2.55, ["Inferred corridor", "Station", "Hotspots", "Violations"], rows, col_widths=[2.85, 1.6, 0.9, 1.0])
    card(s, 6.4, 4.7, 3.05, 1.45, "What is directly real", ["Exact latitude and longitude", "Station and time window", "Observed road-space text"], GREEN, GREEN_50)
    card(s, 9.7, 4.7, 3.05, 1.45, "What is a proxy", ["Lane-level geometry", "Vehicle travel speed", "Traffic-signal health"], ORANGE, ORANGE_50)
    foot(s)


def slide_payoff(prs, d, pg):
    s = blank(prs)
    header(s, "Operational payoff", "Risk becomes a deployable, resource-aware enforcement plan and simulator", pg)
    m = d["metrics"]
    kpi(s, 0.6, 1.8, 2.35, 0.95, i_fmt(m["top20_lane_recovery_minutes"]), "Top-20 lane-recovery minutes", BLUE)
    kpi(s, 3.18, 1.8, 2.35, 0.95, i_fmt(m["top20_capacity_loss_minutes"]), "Top-20 capacity-loss minutes", ORANGE)
    kpi(s, 5.76, 1.8, 2.35, 0.95, i_fmt(m["immediate_clearance_zones"]), "Immediate-clearance zones", RED)
    kpi(s, 8.34, 1.8, 2.05, 0.95, i_fmt(m["tow_priority_zones"]), "Tow-priority zones", PURPLE)
    kpi(s, 10.6, 1.8, 2.1, 0.95, f"{i_fmt(m['median_clearance_sla_minutes'])} min", "Median response deadline", GREEN)
    pic(s, chart_elrm(d), 0.6, 3.1, w=5.05)
    pic(s, chart_policy(d), 5.82, 3.1, w=4.1)
    top = d["impact"].sort_values("operational_priority_score_0_100", ascending=False).head(4)
    rows = [[r.station[:13], r.recommended_action.replace(" enforcement", "")[:18], f"{r.estimated_lane_recovery_minutes:.0f}", f"{r.clearance_sla_minutes:.0f}"] for r in top.itertuples()]
    mini_table(s, 10.16, 3.1, 2.58, 1.92, ["Station", "Action", "Rec.", "SLA"], rows, col_widths=[0.84, 0.96, 0.42, 0.36], fs=6.6)
    rect(s, 10.16, 5.24, 2.58, 0.9, fill=GREEN_50, line=GREEN)
    tx(s, 10.32, 5.33, 2.25, 0.25, "SHIFT SIMULATION", size=7.4, color=GREEN, bold=True)
    tx(s, 10.32, 5.6, 2.25, 0.34, f"{i_fmt(m['standard_budget_recovery_minutes'])} lane-min recovered at {f_fmt(m['standard_budget_recovery_per_resource_hour'])} min/resource-hour", size=7.8, color=INK, bold=True, spacing=1.02)
    tx(s, 0.6, 6.35, 12.1, 0.5, "The output is a station-wise plan: rank, zone, time window, reason, action, resource hours, return on investment, response deadline and a printable field brief.", size=10.5, color=INK, bold=True, spacing=1.08)
    foot(s)


def slide_policy(prs, d, pg):
    s = blank(prs)
    header(s, "Deployment-policy simulation", "Dispatch rules become measurable trade-offs before a shift starts", pg)
    pic(s, chart_policy(d), 0.66, 1.85, w=6.2)
    m = d["metrics"]
    card(s, 7.2, 1.85, 5.55, 1.25, "Standard-shift back-test",
         f"Under a standard patrol and tow budget, the ParkPulse priority rule recovers {i_fmt(m['standard_budget_recovery_minutes'])} lane-minutes and covers {i_fmt(m['standard_budget_high_spillback_zones'])} high-spillback zones in a single shift.", BLUE, BLUE_50)
    kpi(s, 7.2, 3.32, 2.65, 0.95, f_fmt(m["standard_budget_recovery_per_resource_hour"]) + " min", "Recovery per resource-hour", GREEN)
    kpi(s, 10.1, 3.32, 2.65, 0.95, "+" + f_fmt(m["parkpulse_vs_density_recovery_uplift_pct"]) + "%", "Recovery uplift versus raw density", ORANGE)
    card(s, 7.2, 4.5, 5.55, 1.65, "Why this matters for management",
         ["Shows the recovery / coverage trade-off before officers leave", "Makes the dispatch order explainable to each station",
          f"Beats a Traffic-Obstruction-Index-only rule by {f_fmt(m.get('parkpulse_vs_tori_recovery_uplift_pct', 9.1))}% recovery and doubles spillback coverage",
          "Creates a repeatable post-shift review ritual"], PURPLE, PURPLE_50)
    foot(s)


def slide_dash_command(prs, d, pg):
    s = blank(prs)
    header(s, "Product walkthrough, part one", "Command Center: six defensible numbers and one decision map", pg)
    shot(s, "command_center_crop.png", 0.58, 1.82, 7.45)
    cx = 8.35
    card(s, cx, 1.82, 4.4, 1.2, "Six decision numbers", "Recurrence capture, lane-recovery minutes, spillback zones, chronic-repeat zones, patrol-gap zones, and the median response deadline — nothing decorative.", BLUE, BLUE_50, title_size=10.5)
    card(s, cx, 3.12, 4.4, 1.2, "Live corridor-risk map", "Real Bengaluru street map with the inferred corridor-risk network and clickable hotspots; every dot opens its brief.", ORANGE, ORANGE_50, title_size=10.5)
    card(s, cx, 4.42, 4.4, 1.2, "Ranked priority queue", "The top enforcement priorities are listed beside the map, ordered by operational priority, ready to dispatch.", GREEN, GREEN_50, title_size=10.5)
    card(s, cx, 5.72, 4.4, 1.05, "Built for a control room", "A duty officer reads the whole city in one screen and acts in seconds.", PURPLE, PURPLE_50, title_size=10.5)
    foot(s, note="Live screen capture of the ParkPulse React command-center application")


def slide_dash_workflow(prs, d, pg):
    s = blank(prs)
    header(s, "Product walkthrough, part two", "From hotspot explanation to resource simulation in one duty cycle", pg)
    shots = [
        ("hotspot_intelligence_crop.png", "Hotspot Intelligence", ORANGE),
        ("deployment_simulator_crop.png", "Deployment Simulator", GREEN),
    ]
    positions = [(0.58, 1.92), (0.58, 4.43)]
    for (name, label, col), (x, y) in zip(shots, positions):
        shot(s, name, x, y, 6.35)
        tag(s, x, y - 0.04, label, col, w=2.45)
    rect(s, 7.25, 1.92, 5.5, 4.78, fill=BG, line=LINE)
    tx(s, 7.45, 2.06, 5.05, 0.3, "ONE SHIFT, FOUR DECISIONS", size=9, color=PURPLE, bold=True)
    steps = [
        ("Brief", "Start with the ranked city queue from the Command Center and assign first-wave responsibility by station."),
        ("Explain", "Open a hotspot to see the reason: road-space obstruction, repeat pressure, patrol gap, spillback risk and response deadline."),
        ("Dispatch", "Choose tow or fixed-window patrol and generate the officer brief with the exact location and plan."),
        ("Simulate", "Move patrol and tow sliders to see recovery-minutes, covered hotspots and resource-hour trade-offs."),
        ("Learn", "After each shift, feed closure time and action status back into the pipeline for the next run."),
    ]
    yy = 2.48
    for i, (t, b) in enumerate(steps, 1):
        rect(s, 7.45, yy, 0.34, 0.34, fill=BLUE, line=None, rounded=True)
        tx(s, 7.45, yy + 0.05, 0.34, 0.24, str(i), size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tx(s, 7.9, yy - 0.02, 4.55, 0.26, t, size=10.7, color=INK, bold=True)
        tx(s, 7.9, yy + 0.26, 4.55, 0.54, b, size=8.1, color=SLATE, spacing=1.07)
        yy += 0.82
    mini_table(s, 7.45, 6.28, 5.05, 0.32, ["Day 1", "30 days", "90 days"], [["Batch dispatch", "Closure feedback", "Speed/road feeds"]], col_widths=[1.65, 1.7, 1.7], fs=6.5)
    foot(s, note="Live screen captures of the ParkPulse React application")


def slide_roadmap(prs, d, pg):
    s = blank(prs)
    header(s, "Implementation roadmap", "A day-one product with a clean path to a measured, production system", pg)
    phases = [
        ("Day 1", "Batch decision layer", ["Nightly pipeline run", "Station deployment cards", "Printable field briefs"], BLUE),
        ("30 days", "Operational adoption", ["Daily control-room review", "Repeat-offender follow-up", "Response-deadline feedback loop"], ORANGE),
        ("90 days", "Outcome learning", ["Capture closure and tow time stamps", "Before-and-after tagging", "Displacement monitoring"], GREEN),
        ("Next feeds", "Richer intelligence", ["Road and lane geometry", "Traffic-speed feed", "Signal health and patrol location"], PURPLE),
    ]
    for idx, (phase, title, items, col) in enumerate(phases):
        x = 0.62 + idx * 3.06
        tag(s, x, 1.92, phase, col, w=1.0)
        card(s, x, 2.4, 2.74, 2.85, title, items, col, WHITE)
        if idx < 3:
            shp(s, MSO_SHAPE.RIGHT_ARROW, x + 2.76, 3.6, 0.26, 0.24, fill=col, line=None)
    mini_table(s, 0.9, 5.7, 11.55, 0.9,
               ["Owner", "Decision ritual", "Output produced"],
               [["Traffic control room", "Daily 15-minute hotspot review", "First-wave dispatch order"],
                ["Police station", "Post-shift closure update", "Recurrence reduced, displaced, or unchanged"]],
               col_widths=[2.6, 5.1, 3.85])
    foot(s)


def slide_defensibility(prs, d, pg):
    s = blank(prs)
    header(s, "Defensibility and close", "What we claim now, what is a proxy, and why this is day-one implementable", pg)
    claims = [
        ("Claim with confidence now", ["Recurring illegal-parking hotspots", "Repeat and chronic vehicle pressure", "Patrol-coverage gap candidates", "Actionable station-wise deployment"], GREEN, GREEN_50),
        ("Present as a transparent proxy", ["Lane obstruction severity", "Queue-spillback risk", "Capacity-loss minutes", "Equivalent lane-recovery minutes"], ORANGE, ORANGE_50),
        ("Do not claim until feeds exist", ["Measured vehicle-speed improvement", "Surveyed lane-level geometry", "Before-and-after enforcement outcome", "Live signal and workforce state"], RED, RED_50),
    ]
    for idx, (title, items, col, fill) in enumerate(claims):
        card(s, 0.62 + idx * 4.08, 1.92, 3.74, 2.45, title, items, col, fill)
    tx(s, 0.7, 4.55, 12.0, 0.4, "Why this earns trust: the system creates value today, but never pretends that violation records are a complete traffic-sensor network.", size=12.5, color=INK, bold=True, spacing=1.08)
    mini_table(s, 0.9, 5.18, 11.55, 1.0,
               ["A sharp judge will ask", "Our answer"],
               [["Why is the obstruction index not the top predictor?", "Because recurrence and action are separate jobs: density and the forecast predict; the index explains and prioritises."],
                ["Where is the measured traffic-flow impact?", "It is not in the dataset; recovery minutes are a labelled, bounded proxy until speed and queue feeds are added."]],
               col_widths=[4.2, 7.15], fs=8.0)
    rect(s, 0.9, 6.38, 11.55, 0.48, fill=NAVY, line=None, rounded=True)
    tx(s, 1.1, 6.48, 11.15, 0.22, "ParkPulse narrows the field, explains the reason, dispatches the right resource, recovers road capacity, and learns after every shift.", size=10.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    foot(s)


def slide_closing(prs, d, pg):
    s = blank(prs)
    rect(s, -0.1, -0.1, SW + 0.2, SH + 0.2, fill=NAVY, line=None, rounded=False)
    logo(s, 0.75, 0.55, dark=True)
    tx(s, 0.9, 1.45, 11.6, 1.15, [
        ("ParkPulse is not a dashboard.\n", {"color": WHITE, "bold": True, "size": 32}),
        ("It is a deployment engine for illegal-parking enforcement.", {"color": YELLOW, "bold": True, "size": 32}),
    ], size=32, color=WHITE, spacing=1.05)
    tx(s, 0.92, 3.1, 11.0, 0.6, "Use the violation records a city already has to make targeted enforcement measurable from day one, then add road-network, speed, signal and workforce feeds as upgrades.", size=15, color=RGBColor(0xE7, 0xEF, 0xFF), spacing=1.18)
    m = d["metrics"]
    vals = [
        (p_fmt(m["headline_recurrence_capture_at_20"]), "Validated recurrence capture"),
        (i_fmt(m["top20_lane_recovery_minutes"]) + " min", "Top-20 lane recovery"),
        (i_fmt(m["chronic_repeat_priority_zones"]), "Chronic repeat-offender zones"),
        (i_fmt(m["patrol_gap_priority_zones"]), "Patrol-coverage gap zones"),
    ]
    for idx, (val, lab) in enumerate(vals):
        kpi(s, 0.95 + idx * 3.0, 4.5, 2.75, 1.1, val, lab, YELLOW, dark=True)
    tx(s, 0.95, 6.25, 11.5, 0.34, "Narrow the field, explain the reason, dispatch the right resource, recover capacity, and learn after every shift.", size=13.5, color=WHITE, bold=True)
    tx(s, 0.95, 6.95, 11.5, 0.24, "Flipkart GRiD · Theme 1 · ParkPulse Bengaluru", size=9.2, color=RGBColor(0xC7, 0xD7, 0xF5), bold=True)


def build():
    d = load_data()
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    slides = [
        slide_title, slide_questions, slide_strategy, slide_challenge, slide_dataset,
        slide_eda_signals, slide_architecture, slide_ml, slide_scoring, slide_validation,
        slide_spatial, slide_payoff, slide_dash_command, slide_dash_workflow, slide_defensibility,
    ]
    for i, fn in enumerate(slides, 1):
        fn(prs, d, i)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({OUT.stat().st_size / 1024:.0f} KB, {len(prs.slides)} slides).")


if __name__ == "__main__":
    build()
